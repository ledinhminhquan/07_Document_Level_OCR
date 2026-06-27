"""Generate the submission slides.pptx (python-pptx) — ~10 concise visual slides.
Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, artifacts_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import load_artifacts

logger = get_logger(__name__)


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    ev = arts.get("eval") or {}
    m = ev.get("model", {})
    before, after, red = m.get("cer_before"), m.get("cer_after"), m.get("cer_reduction_rel")
    res = (f"CER {before:.3f} -> {after:.3f} ({red*100:.0f}% reduction)" if (before is not None and after is not None and red is not None)
           else "train + evaluate to populate results")
    return [
        ("Document-Level OCR System",
         [f"{cfg.author} — Student {cfg.student_id}", "NLP in Industry — Final Assignment",
          "Document image/PDF -> clean, structured, ordered text", "OCR + layout + a trainable post-OCR corrector + an agent"]),
        ("Business Problem & Motivation",
         ["Flat OCR = unordered, error-riddled text dump", "Wrong reading order on multi-column pages; no structure",
          "Char errors: rn->m, 0<->O, 1<->l, merged/split words", "Clean structured text feeds search, RAG, analytics"]),
        ("Proposed NLP Solution",
         ["Trainable core: ByT5 post-OCR corrector (char-level)", "Engine-agnostic uplift on top of any OCR backend",
          "Layout + reading order -> Markdown / JSON structure", "An agent orchestrates document -> structured text"]),
        ("System Architecture",
         ["page/PDF -> preprocess -> born-digital/scanned route", "-> layout + reading order -> OCR regions",
          "-> POST-OCR CORRECT (trained) -> assemble", "-> Markdown + JSON blocks + plain text"]),
        ("Data Overview",
         ["Primary = synthetic OCR-noise generator (rn->m, 0<->O ...)", "Real mix: PleIAs/Post-OCR-Correction (english, CC0)",
          "Identity baseline CER ~0.088 / WER ~0.49 to beat", "ICDAR-2019 as an optional benchmark"]),
        ("Model & Evaluation Results",
         [res, "Headline = CER/WER REDUCTION vs raw OCR",
          "Safety gate: % improved must ≫ % degraded", "Beats identity + SymSpell dictionary baselines"]),
        ("Agentic AI Component",
         ["Deterministic FSM + optional LLM brain (rule fallback)", "D1 page-quality · D2 born-digital/scanned routing",
          "D3 OCR-confidence gate · D4 correction acceptance", "Edit-budget guard so corrections can't hallucinate"]),
        ("Deployment Overview",
         ["FastAPI /ocr (image/PDF) + /correct (text)", "Gradio UI · CLI batch · Docker · HF Space",
          "Born-digital pages skip OCR (D2) -> ~200 ms/page", "CPU-able default (Tesseract + ByT5/t5-small)"]),
        ("Ethics, Privacy & Risks",
         ["Documents may contain PII -> minimize retention/logging", "Corrector must not 'correct' away meaning (edit budget)",
          "OCR of private/copyright docs needs rights/consent", "Sim-to-real gap -> report real PleIAs CER too"]),
        ("Key Takeaways & Future Work",
         ["A small trained corrector gives engine-agnostic CER uplift", "Synthetic noise unlocks training without paired scans",
          "Future: table/figure structure, layout model, multi-lingual", "Future: confidence calibration, active learning"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else artifacts_dir() / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    chart = charts_mod.cer_reduction_chart(arts.get("eval") or {}, out_path.parent / "charts" / "slide_cer.png")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(30); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.3 if (i == 5 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "•  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 5 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} — {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
